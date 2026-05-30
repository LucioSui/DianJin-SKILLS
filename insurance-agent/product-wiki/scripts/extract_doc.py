#!/usr/bin/env python3
"""
保险代理产品 Wiki - 文档提取脚本
支持格式：PDF、Word(.docx)、PPT(.pptx)、纯文本(.txt/.md)
用法: python3 extract_doc.py <file_path> [--output output.md]
"""

import sys
import os
import argparse
from pathlib import Path
from datetime import datetime


def extract_pdf(file_path: str) -> str:
    """提取 PDF 文本内容"""
    try:
        import pdfplumber
    except ImportError:
        print("错误: 需要安装 pdfplumber", file=sys.stderr)
        print("运行: pip install pdfplumber", file=sys.stderr)
        sys.exit(1)

    text_parts = []
    with pdfplumber.open(file_path) as pdf:
        total_pages = len(pdf.pages)
        print(f"PDF 共 {total_pages} 页，正在提取...", file=sys.stderr)
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text:
                text_parts.append(f"\n--- 第 {i} 页 ---\n{text}")
            if i % 10 == 0:
                print(f"  已处理 {i}/{total_pages} 页", file=sys.stderr)
    return "\n".join(text_parts)


def extract_docx(file_path: str) -> str:
    """提取 Word 文档内容"""
    try:
        from docx import Document
    except ImportError:
        print("错误: 需要安装 python-docx", file=sys.stderr)
        print("运行: pip install python-docx", file=sys.stderr)
        sys.exit(1)

    doc = Document(file_path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    # 提取表格
    tables = doc.tables
    if tables:
        parts.append("\n## 文档中的表格\n")
        for i, table in enumerate(tables, 1):
            parts.append(f"\n### 表格 {i}\n")
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                parts.append("| " + " | ".join(cells) + " |")
    return "\n".join(parts)


def extract_pptx(file_path: str) -> str:
    """提取 PPT 演示文稿内容"""
    try:
        from pptx import Presentation
    except ImportError:
        print("错误: 需要安装 python-pptx", file=sys.stderr)
        print("运行: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"\n--- 第 {i} 张幻灯片 ---\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
            if shape.has_table:
                parts.append(f"\n[表格内容]\n")
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append("| " + " | ".join(cells) + " |")
    return "\n".join(parts)


def extract_text(file_path: str) -> str:
    """提取纯文本文件"""
    encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, LookupError):
            continue
    print(f"错误: 无法解码文件 {file_path}", file=sys.stderr)
    sys.exit(1)


def extract_image(file_path: str) -> str:
    """图片提取 - 返回提示信息，由 Agent 使用 vision_analyze 处理"""
    print("[提示] 图片 OCR 需由 Agent 调用 vision_analyze 工具处理", file=sys.stderr)
    print(f"[提示] 图片路径: {file_path}", file=sys.stderr)
    return f"[图片文件: {os.path.basename(file_path)} - 请 Agent 使用 vision_analyze 提取内容]"


def main():
    parser = argparse.ArgumentParser(description="从文档中提取文本内容")
    parser.add_argument("file", help="输入文件路径")
    parser.add_argument("--output", "-o", help="输出文件路径（默认输出到 stdout）")
    args = parser.parse_args()

    file_path = args.file
    if not os.path.exists(file_path):
        print(f"错误: 文件不存在: {file_path}", file=sys.stderr)
        sys.exit(1)

    ext = Path(file_path).suffix.lower()
    extractors = {
        ".pdf": extract_pdf,
        ".docx": extract_docx,
        ".doc": extract_docx,
        ".pptx": extract_pptx,
        ".ppt": extract_pptx,
        ".txt": extract_text,
        ".md": extract_text,
        ".text": extract_text,
        ".png": extract_image,
        ".jpg": extract_image,
        ".jpeg": extract_image,
        ".gif": extract_image,
        ".bmp": extract_image,
        ".webp": extract_image,
    }

    extractor = extractors.get(ext)
    if extractor is None:
        print(f"错误: 不支持的文件格式: {ext}", file=sys.stderr)
        print(f"支持的格式: {', '.join(extractors.keys())}", file=sys.stderr)
        sys.exit(1)

    print(f"正在提取 {ext} 文件: {file_path}", file=sys.stderr)
    content = extractor(file_path)

    # 添加元数据头
    header = (
        f"# 文档提取内容\n\n"
        f"> 源文件: {os.path.basename(file_path)}\n"
        f"> 提取时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}\n"
        f"> 格式: {ext}\n"
        f"> 字符数: {len(content):,}\n\n"
        f"---\n\n"
    )

    result = header + content

    if args.output:
        with open(args.output, "w", encoding="utf-8") as f:
            f.write(result)
        print(f"已保存到: {args.output}", file=sys.stderr)
    else:
        print(result)


if __name__ == "__main__":
    main()
